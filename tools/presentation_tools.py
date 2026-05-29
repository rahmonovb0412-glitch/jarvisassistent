"""
Prezentatsiya yaratish toollari - python-pptx asosida
Professional PowerPoint (.pptx) slaydlar yaratadi
"""

import os
import json

WORKSPACE = os.getenv("WORKSPACE_DIR", "./workspace")

# Tayyor dizayn mavzulari (fon + matn ranglari)
THEMES = {
    "dark":    {"bg": (10, 10, 25),    "title": (79, 142, 247),  "text": (220, 230, 245)},
    "light":   {"bg": (248, 250, 252),  "title": (30, 58, 138),   "text": (30, 41, 59)},
    "ocean":   {"bg": (8, 47, 73),      "title": (56, 189, 248),  "text": (224, 242, 254)},
    "sunset":  {"bg": (40, 12, 30),     "title": (251, 146, 60),  "text": (254, 235, 220)},
    "emerald": {"bg": (6, 38, 30),      "title": (52, 211, 153),  "text": (220, 252, 240)},
}



def create_presentation(title: str, slides, output_filename: str = "presentation.pptx", theme: str = "dark") -> dict:
    """
    Professional prezentatsiya (.pptx) yaratadi.
    title: prezentatsiya sarlavhasi (birinchi slayd)
    slides: slaydlar ro'yxati. JSON yoki list:
            [{"heading": "Sarlavha", "points": ["fikr1", "fikr2"]}]
            yoki [{"heading": "...", "text": "..."}]
    theme: dark, light, ocean, sunset, emerald
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return {"success": False, "error": "python-pptx o'rnatilmagan: pip install python-pptx"}

    try:
        # slides string bo'lsa JSON parse qilish
        if isinstance(slides, str):
            try:
                slides = json.loads(slides)
            except json.JSONDecodeError:
                slides = [{"heading": title, "points": [slides]}]
        if not isinstance(slides, list):
            slides = [slides]

        th = THEMES.get(theme, THEMES["dark"])
        bg = RGBColor(*th["bg"])
        c_title = RGBColor(*th["title"])
        c_text = RGBColor(*th["text"])

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        def _fill_bg(slide):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = bg

        # ── Sarlavha slaydi ──
        s = prs.slides.add_slide(blank)
        _fill_bg(s)
        tb = s.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(2))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title
        r.font.size = Pt(48); r.font.bold = True; r.font.color.rgb = c_title

        # ── Mazmun slaydlari ──
        for sl in slides:
            s = prs.slides.add_slide(blank)
            _fill_bg(s)
            heading = sl.get("heading") or sl.get("title") or ""
            # Sarlavha
            ht = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
            hp = ht.text_frame.paragraphs[0]
            hr = hp.add_run(); hr.text = heading
            hr.font.size = Pt(34); hr.font.bold = True; hr.font.color.rgb = c_title
            # Mazmun
            body = s.shapes.add_textbox(Inches(1), Inches(1.9), Inches(11.3), Inches(5))
            btf = body.text_frame; btf.word_wrap = True
            points = sl.get("points")
            if not points:
                txt = sl.get("text", "")
                points = [txt] if txt else []
            for i, pt in enumerate(points):
                para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                pr = para.add_run(); pr.text = "•  " + str(pt)
                pr.font.size = Pt(22); pr.font.color.rgb = c_text
                para.space_after = Pt(14)

        os.makedirs(WORKSPACE, exist_ok=True)
        out = os.path.join(WORKSPACE, output_filename) if not os.path.isabs(output_filename) else output_filename
        prs.save(out)

        return {
            "success": True,
            "message": f"✅ Prezentatsiya tayyor: {output_filename}",
            "path": out,
            "slides": len(prs.slides),
            "theme": theme,
        }
    except Exception as e:
        return {"success": False, "error": str(e)}


def list_presentations() -> dict:
    """Yaratilgan prezentatsiyalar ro'yxati"""
    try:
        items = []
        for f in os.listdir(WORKSPACE):
            if f.lower().endswith((".pptx", ".ppt")):
                p = os.path.join(WORKSPACE, f)
                items.append({"name": f, "size": f"{os.path.getsize(p)//1024} KB"})
        return {"success": True, "count": len(items), "presentations": items}
    except Exception as e:
        return {"success": False, "error": str(e)}
